from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy

W = Inches(13.33)
H = Inches(7.5)

def rgb(r, g, b):
    return RGBColor(r, g, b)

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, font_size=14, bold=False, color=rgb(0,0,0),
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_table(slide, data, headers, l, t, w, h,
              hdr_bg=rgb(11,37,69), hdr_fg=rgb(255,255,255),
              row_bg1=rgb(255,255,255), row_bg2=rgb(244,247,252),
              font_size=11):
    rows = len(data) + 1
    cols = len(headers)
    table = slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h)).table
    col_w = Inches(w / cols)
    for i in range(cols):
        table.columns[i].width = col_w
    # Header
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hdr_bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = hdr
        run.font.bold = True
        run.font.size = Pt(font_size)
        run.font.color.rgb = hdr_fg
    # Data rows
    for ri, row in enumerate(data):
        bg = row_bg1 if ri % 2 == 0 else row_bg2
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(font_size)
            run.font.color.rgb = rgb(60, 60, 60)
    return table

def slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

# ─────────────────────────────────────────────
# MINISTER / MP DECK
# ─────────────────────────────────────────────
def build_minister_deck():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    NAVY   = rgb(11, 37, 69)
    BLUE   = rgb(27, 79, 138)
    GREEN  = rgb(19, 136, 8)
    SAFFRON= rgb(255, 153, 51)
    WHITE  = rgb(255, 255, 255)
    LGRAY  = rgb(244, 247, 252)
    GOLD   = rgb(255, 209, 0)

    blank = prs.slide_layouts[6]

    # ── SLIDE 1: COVER ────────────────────────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, NAVY)
    # Tricolor stripe
    add_rect(sl, 0, 0, 4.44, 0.12, fill_rgb=SAFFRON)
    add_rect(sl, 4.44, 0, 4.44, 0.12, fill_rgb=WHITE)
    add_rect(sl, 8.88, 0, 4.45, 0.12, fill_rgb=GREEN)
    add_text(sl, "DigiSampatti", 0, 1.2, 13.33, 1.4, font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, "ಡಿಜಿ ಸಂಪತ್ತಿ", 0, 2.5, 13.33, 0.7, font_size=26, bold=False, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(sl, "Karnataka's Property Legal Intelligence Platform", 0, 3.1, 13.33, 0.6, font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, "Making Government Data Understandable for Every Citizen", 0, 3.65, 13.33, 0.5, font_size=16, color=rgb(200,210,230), align=PP_ALIGN.CENTER)
    add_text(sl, "Presented to the Hon'ble Revenue Minister / MP, Government of Karnataka", 0, 4.3, 13.33, 0.5, font_size=14, color=rgb(160,180,210), align=PP_ALIGN.CENTER, italic=True)
    # Badges
    for i, badge in enumerate(["Patent Pending", "Submitted to Google Play Store", "DPDP Act 2023 Compliant", "Digital India Aligned"]):
        add_rect(sl, 0.8 + i*3.1, 5.3, 2.8, 0.45, fill_rgb=rgb(30,55,100), line_rgb=rgb(100,140,200), line_width=1)
        add_text(sl, badge, 0.82 + i*3.1, 5.33, 2.76, 0.42, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 0, 7.38, 4.44, 0.12, fill_rgb=SAFFRON)
    add_rect(sl, 4.44, 7.38, 4.44, 0.12, fill_rgb=WHITE)
    add_rect(sl, 8.88, 7.38, 4.45, 0.12, fill_rgb=GREEN)

    # ── SLIDE 2: THE PROBLEM ──────────────────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, WHITE)
    add_rect(sl, 0, 0, 13.33, 0.08, fill_rgb=NAVY)
    add_text(sl, "The Problem", 0.5, 0.15, 12, 0.4, font_size=11, bold=True, color=BLUE)
    add_text(sl, "What Is Happening to Karnataka Citizens Today", 0.5, 0.5, 12.33, 0.8, font_size=28, bold=True, color=NAVY)
    add_text(sl, "One search. Seven portals. Ten minutes. Done from home. No middlemen.", 0.5, 1.2, 12, 0.45, font_size=15, italic=True, color=rgb(120,120,120))

    stories = [
        ("Raju — Farmer, Tumkur",
         "Paid ₹40 lakhs — family life savings — on 2 acres. Six months later: court summons. The land had 3 prior encumbrances in Bhoomi. He never checked. He did not know how. He lost everything."),
        ("Savitha — Govt Teacher, Mysuru",
         "Bought ₹32 lakh flat with PF savings. Builder said 'all approvals done.' Two years later: no RERA registration, active court dispute. Bank filed recovery. She pays EMI on a flat she may lose."),
        ("Suresh — NRI Engineer, Dubai",
         "Sent ₹22 lakhs for a plot in Dharwad. Could not visit India. Trusted the broker. Title was forged — visible in IGRS records. Could not check from UAE. Lost everything. Case still in court."),
    ]
    for i, (name, story) in enumerate(stories):
        add_rect(sl, 0.5, 1.75 + i*1.45, 12.33, 1.3, fill_rgb=rgb(255,253,231), line_rgb=rgb(255,224,130), line_width=1)
        add_text(sl, name, 0.7, 1.8 + i*1.45, 12, 0.35, font_size=13, bold=True, color=NAVY)
        add_text(sl, story, 0.7, 2.1 + i*1.45, 12, 0.8, font_size=12, color=rgb(80,80,80))

    add_rect(sl, 0.5, 6.15, 12.33, 0.75, fill_rgb=rgb(255,235,238), line_rgb=rgb(198,40,40), line_width=1)
    add_text(sl, "Karnataka: 1.2 crore property transactions/year.  90% of buyers verify NOTHING before paying.  The government has the data. The citizen cannot reach it.", 0.65, 6.2, 12, 0.65, font_size=13, bold=True, color=rgb(123,26,26))

    # ── SLIDE 3: WHAT DIGISAMPATTI DOES ───────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, WHITE)
    add_rect(sl, 0, 0, 13.33, 0.08, fill_rgb=NAVY)
    add_text(sl, "The Solution", 0.5, 0.15, 12, 0.4, font_size=11, bold=True, color=BLUE)
    add_text(sl, "What DigiSampatti Does — 5 Steps, From Home", 0.5, 0.5, 12.33, 0.8, font_size=28, bold=True, color=NAVY)

    steps = [
        ("1", "No Login to 7 Portals — One Search",
         "Enter survey number OR scan document photo OR stand on land (GPS). DigiSampatti handles all 7 portals automatically. No visiting any government website."),
        ("2", "Documents Retrieved Simultaneously",
         "All 7 portal records retrieved at the same moment for one property. This alone is impossible to do manually — by the time you collect document 7, document 1 may have changed."),
        ("3", "Timestamped Proof — Stored Securely",
         "Encrypted, timestamped snapshot created instantly. This is your legal proof of due diligence — court-admissible. Stored 30 days on our servers, yours forever on your phone."),
        ("4", "Cross-Portal Contradiction Detection",
         "All 7 documents read against each other. Fraud lives in the gap between portals — not inside any one portal. 47 Karnataka land law rules applied automatically."),
        ("5", "One Verdict in Plain Kannada — SAFE / CAUTION / DO NOT BUY",
         "The farmer in Tumkur does not need 7 PDFs. He needs one decision. DigiSampatti gives him that in 10 minutes for ₹99 — before paying ₹7 lakh in stamp duty."),
    ]
    for i, (num, title, desc) in enumerate(steps):
        y = 1.35 + i * 1.1
        add_rect(sl, 0.5, y, 0.55, 0.9, fill_rgb=BLUE)
        add_text(sl, num, 0.5, y + 0.1, 0.55, 0.7, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, 1.15, y, 11.68, 0.9, fill_rgb=LGRAY, line_rgb=rgb(208,220,240), line_width=1)
        add_text(sl, title, 1.25, y + 0.02, 11.4, 0.35, font_size=13, bold=True, color=NAVY)
        add_text(sl, desc, 1.25, y + 0.38, 11.4, 0.48, font_size=11, color=rgb(80,80,80))

    # ── SLIDE 4: 7 PORTALS ────────────────────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, WHITE)
    add_rect(sl, 0, 0, 13.33, 0.08, fill_rgb=NAVY)
    add_text(sl, "The Ecosystem", 0.5, 0.15, 12, 0.4, font_size=11, bold=True, color=BLUE)
    add_text(sl, "7 Government Systems — Connected for the First Time", 0.5, 0.5, 12.33, 0.7, font_size=28, bold=True, color=NAVY)
    add_text(sl, "No citizen can check all 7 manually. No government portal links them. DigiSampatti cross-references all 7 simultaneously.", 0.5, 1.15, 12.33, 0.45, font_size=13, italic=True, color=rgb(100,100,100))

    headers = ["Portal", "What We Check", "What Citizens Miss Without DigiSampatti"]
    data = [
        ["Bhoomi (Revenue Dept)", "Title, mutations, govt acquisition, PTCL restrictions", "Column 11/12 meaning, SC/ST violations, incomplete mutations"],
        ["IGRS / KAVERI", "Encumbrance Certificate — all loans and charges", "Hidden mortgages, uncleared loans, multiple sale deeds"],
        ["RERA Karnataka", "Builder registration, project approval, complaints", "Unregistered projects, RERA vs Bhoomi land-type contradiction"],
        ["BBMP / BDA / CMC", "Khata status (A/B/E), tax dues, layout approval", "B-Khata trap, unapproved layouts, tax arrears blocking mutation"],
        ["eCourts (National)", "Court cases — reverse lookup by survey number + owner name", "Active disputes citizens NEVER find (eCourts needs case number, not survey number)"],
        ["CERSAI (Central)", "All bank charges and mortgages nationwide", "Bank loans not visible in state IGRS, NRI property loans"],
        ["Benami Portal (Income Tax)", "Properties under Benami investigation", "Govt can SEIZE property even after purchase. Most buyers never check this."],
    ]
    add_table(sl, data, headers, 0.5, 1.65, 12.33, 5.2, font_size=11)

    # ── SLIDE 5: FULL JOURNEY ─────────────────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, WHITE)
    add_rect(sl, 0, 0, 13.33, 0.08, fill_rgb=NAVY)
    add_text(sl, "Full Journey", 0.5, 0.15, 12, 0.4, font_size=11, bold=True, color=BLUE)
    add_text(sl, "From Advance Payment to Final Mutation — All From Home", 0.5, 0.5, 12.33, 0.7, font_size=27, bold=True, color=NAVY)

    headers = ["Journey Step", "Problem Today (Without DigiSampatti)", "DigiSampatti Solution (From Home)"]
    data = [
        ["Before paying advance", "Pays ₹1–5L token money blindly. Discovers issues after paying.", "Pre-advance check: full 7-portal report before any rupee changes hands."],
        ["Property identification", "Doesn't know survey number. Can't navigate portal.", "Camera scans document. GPS identifies land on-site. Auto-extracted."],
        ["Full verification", "3–5 days, 7 portals, ₹15,000 lawyer, English docs, no verdict.", "10 min, ₹99, 7 portals, plain Kannada verdict. Govt fees paid."],
        ["Legal opinion", "₹50,000 lawyer. Rural citizens have no access.", "Verified Karnataka lawyer on app. ₹500–2,000. 2 hours. Remote."],
        ["Stamp duty calculation", "Surprise at sub-registrar office. Transaction fails.", "Exact calculation by property type, value, district — upfront."],
        ["Post-registration mutation", "No notification. Manual follow-up for months.", "Auto-SMS when Bhoomi mutation confirmed. Legal proof sent."],
    ]
    add_table(sl, data, headers, 0.5, 1.3, 12.33, 5.6, font_size=11)

    # ── SLIDE 6: INTELLIGENCE LAYER ───────────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, WHITE)
    add_rect(sl, 0, 0, 13.33, 0.08, fill_rgb=NAVY)
    add_text(sl, "Why Intelligence Matters", 0.5, 0.15, 12, 0.4, font_size=11, bold=True, color=BLUE)
    add_text(sl, "Government Portals Have the Data. Fraud Lives in the Gap Between Them.", 0.5, 0.5, 12.33, 0.8, font_size=26, bold=True, color=NAVY)

    add_rect(sl, 0.5, 1.35, 12.33, 0.6, fill_rgb=rgb(255,235,238), line_rgb=rgb(198,40,40), line_width=1)
    add_text(sl, '"Sir, government portals are all correct. Fraud is never in one portal. Fraud lives in the gap between portals. DigiSampatti reads the gap."', 0.65, 1.38, 12, 0.54, font_size=13, bold=True, color=rgb(123,26,26))

    headers = ["What Portal Shows", "What Citizen Understands", "What DigiSampatti Tells Them"]
    data = [
        ["RERA: Project complete. Bhoomi: Land is agricultural. BBMP: No khata.",
         "Each portal looks fine. Citizen sees no problem.",
         "FRAUD DETECTED: Building is on unconverted agricultural land. Apartment legally invalid. DO NOT BUY."],
        ["IGRS: SBI charge ₹18L — 2022 — ACTIVE",
         "Sees a bank name. Doesn't know if cleared.",
         "CAUTION: Active SBI mortgage. Verify NOC before paying. Cannot get clear title."],
        ["eCourts: Case #432/2021 — Tumkur Civil — Active",
         "Cannot find. eCourts needs case number, not survey number.",
         "Active court dispute found on this property. DO NOT BUY. Next hearing: April 2026."],
        ["Benami portal: Survey 45/3 flagged",
         "Has never heard of Benami portal. Buys property.",
         "CRITICAL: Govt may seize this property. DO NOT BUY under any circumstances."],
    ]
    add_table(sl, data, headers, 0.5, 2.05, 12.33, 4.9, font_size=11)

    # ── SLIDE 7: EXPERT COUNCIL ───────────────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, WHITE)
    add_rect(sl, 0, 0, 13.33, 0.08, fill_rgb=NAVY)
    add_text(sl, "Expert Council + Government Fees", 0.5, 0.15, 12, 0.4, font_size=11, bold=True, color=BLUE)
    add_text(sl, "Government Revenue Protected. Retired Karnataka Officials Earning.", 0.5, 0.5, 12.33, 0.7, font_size=27, bold=True, color=NAVY)

    add_rect(sl, 0.5, 1.28, 6.2, 3.5, fill_rgb=LGRAY, line_rgb=rgb(200,215,240), line_width=1)
    add_text(sl, "Government Portal Fees — Fully Paid", 0.7, 1.35, 5.8, 0.4, font_size=14, bold=True, color=NAVY)
    fees_text = ("DigiSampatti pays every government fee on the citizen's behalf:\n\n"
                 "→ Bhoomi RTC:  ₹15–25 paid to Revenue Dept\n"
                 "→ IGRS/KAVERI EC:  ₹30–50 paid to IGRS\n"
                 "→ CERSAI:  ₹10 paid to Central Registry\n"
                 "→ BBMP:  ₹100–500 paid to BBMP\n\n"
                 "Government revenue INCREASES — citizens who\n"
                 "would never have accessed portals individually\n"
                 "now pay fees through DigiSampatti.")
    add_text(sl, fees_text, 0.7, 1.75, 5.8, 2.9, font_size=12, color=rgb(60,60,60))

    add_rect(sl, 7.0, 1.28, 5.83, 3.5, fill_rgb=LGRAY, line_rgb=rgb(200,215,240), line_width=1)
    add_text(sl, "Karnataka Property Intelligence Council", 7.2, 1.35, 5.4, 0.4, font_size=14, bold=True, color=NAVY)
    council_text = ("Just as Anthropic hired 20 financial analysts\nto train Claude — DigiSampatti builds the\nsame model for Karnataka property law:\n\n"
                    "→ Retired Revenue Inspectors (5)\n   Train AI on Bhoomi/mutation complexities\n\n"
                    "→ Bar Council Lawyers (8)\n   Review CAUTION verdicts and land law\n\n"
                    "→ Bank Loan Officers (5)\n   Train AI on CERSAI/mortgage patterns\n\n"
                    "→ Retired Sub-registrar Officers (2)\n   Train AI on registration red flags\n\n"
                    "Each expert earns ₹200–500 per case review.")
    add_text(sl, council_text, 7.2, 1.75, 5.5, 2.9, font_size=12, color=rgb(60,60,60))

    add_rect(sl, 0.5, 4.95, 12.33, 0.9, fill_rgb=rgb(232,245,233), line_rgb=GREEN, line_width=2)
    add_text(sl, "Every DigiSampatti check: Govt portal fees paid  +  Retired Karnataka officials earning  +  Citizen protected\nThe government's infrastructure earns more. The government's retired professionals earn more. Citizens are safer.", 0.65, 4.98, 12, 0.84, font_size=13, bold=True, color=rgb(27,69,32))

    # ── SLIDE 8: HOW IT HELPS GOVERNMENT ──────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, WHITE)
    add_rect(sl, 0, 0, 13.33, 0.08, fill_rgb=NAVY)
    add_text(sl, "Government Benefits", 0.5, 0.15, 12, 0.4, font_size=11, bold=True, color=BLUE)
    add_text(sl, "How DigiSampatti Strengthens the Government", 0.5, 0.5, 12.33, 0.7, font_size=28, bold=True, color=NAVY)

    cards = [
        ("Protects ₹7L Stamp Duty Revenue", "A citizen paying ₹7 lakh in stamp duty deserves to know for ₹99 whether that property is worth registering. DigiSampatti catches issues BEFORE registration. Fraudulent registrations prevented. Stamp duty revenue goes to clean transactions only."),
        ("Fewer Cases in Revenue Courts", "DigiSampatti catches disputes before the transaction. Fewer fraudulent transactions = fewer Revenue Tribunal and Civil Court cases. Karnataka courts are backlogged with preventable property disputes."),
        ("SC/ST Land Protected", "Karnataka's PTCL Act restricts SC/ST land transfer. DigiSampatti flags this in every report — enforcing the government's constitutional obligation at the transaction level. Every check protects the most vulnerable landowners."),
        ("₹137Cr Investment Reaches Every Citizen", "Govt spent ₹137Cr digitising land records. Today 90% of citizens can't use that data. DigiSampatti makes that investment reach the farmer in Tumkur — the person it was meant to serve."),
    ]
    for i, (title, desc) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.45
        y = 1.3 + row * 2.6
        add_rect(sl, x, y, 6.2, 2.4, fill_rgb=LGRAY, line_rgb=rgb(200,215,240), line_width=1)
        add_text(sl, title, x+0.15, y+0.1, 5.9, 0.45, font_size=14, bold=True, color=NAVY)
        add_text(sl, desc, x+0.15, y+0.52, 5.9, 1.75, font_size=12, color=rgb(70,70,70))

    # ── SLIDE 9: THREE REQUESTS ───────────────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, WHITE)
    add_rect(sl, 0, 0, 13.33, 0.08, fill_rgb=NAVY)
    add_text(sl, "Our Requests", 0.5, 0.15, 12, 0.4, font_size=11, bold=True, color=BLUE)
    add_text(sl, "Three Requests — No Government Budget Required", 0.5, 0.5, 12.33, 0.7, font_size=27, bold=True, color=NAVY)
    add_text(sl, "We are not asking for funding. We are asking for permission to serve.", 0.5, 1.15, 12, 0.4, font_size=14, italic=True, color=rgb(100,100,100))

    requests = [
        ("Request 1 — Official API Access",
         "Dedicated, stable API access to Bhoomi, IGRS/KAVERI, and RERA Karnataka for authorized platforms like DigiSampatti.\n\nCurrently accessing public data — which breaks when portal UI changes. A stable API:\n→ Makes reports 10x faster and more reliable\n→ Zero additional load on public portals\n→ Can be conditional on DPDP Act compliance audit\n→ Formal request already submitted"),
        ("Request 2 — eCourts Survey Number Search",
         "Request to NIC/eCourts team to allow property survey number as a search parameter in ecourts.gov.in.\n\nCurrently eCourts only searches by case number — which no buyer has. This single change:\n→ Benefits every Indian property buyer — not just DigiSampatti users\n→ Closes the single largest gap in property due diligence in India\n→ A national public good — minor technical change, enormous citizen benefit"),
        ("Request 3 — Sub-Registrar Pilot + Endorsement",
         "Permission to run DigiSampatti at 5–10 sub-registrar offices in Karnataka as a voluntary citizen service.\n\nA QR code at the counter. Purely voluntary. No mandatory requirement.\n\nOptional: A letter of endorsement from the Revenue Department stating DigiSampatti is a recognized citizen tool. This one letter builds instant trust with farmers who trust government more than any private company."),
    ]
    for i, (title, desc) in enumerate(requests):
        x = 0.4 + i * 4.25
        add_rect(sl, x, 1.65, 4.05, 5.0, fill_rgb=LGRAY, line_rgb=BLUE, line_width=2)
        add_text(sl, title, x+0.12, 1.72, 3.82, 0.5, font_size=13, bold=True, color=NAVY)
        add_text(sl, desc, x+0.12, 2.22, 3.82, 4.3, font_size=11, color=rgb(60,60,60))

    # ── SLIDE 10: CLOSING ─────────────────────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, NAVY)
    add_rect(sl, 0, 0, 4.44, 0.12, fill_rgb=SAFFRON)
    add_rect(sl, 4.44, 0, 4.44, 0.12, fill_rgb=WHITE)
    add_rect(sl, 8.88, 0, 4.45, 0.12, fill_rgb=GREEN)
    add_text(sl, "One Request, Sir", 0, 1.1, 13.33, 0.9, font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    closing = ("The Government of Karnataka has already done the hardest work —\n"
               "digitizing every land record, building Bhoomi, KAVERI, RERA, eCourts.\n\n"
               "The farmer in Tumkur. The teacher in Mysuru. The NRI in Dubai. The SC/ST widow in Kolar.\n"
               "They cannot access or understand that data today — sitting at home, without a middleman, without an officer.\n\n"
               "DigiSampatti bridges that gap — connecting 7 government portals,\n"
               "reading the contradictions between them, and delivering one clear decision in Kannada.\n\n"
               "We are not asking for money.\n"
               "We are asking for the right to make your data work for every citizen.")
    add_text(sl, closing, 1.0, 1.9, 11.33, 4.2, font_size=15, color=rgb(220,230,245), align=PP_ALIGN.CENTER)
    add_text(sl, "₹99. 10 minutes. Before paying ₹7 lakh to your government. That is DigiSampatti.", 1.0, 5.85, 11.33, 0.45, font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(sl, "Sanjay R  |  rajsanjay381@gmail.com  |  +91 89043 42255  |  Patent Pending  |  DPDP Act 2023 Compliant", 0, 6.6, 13.33, 0.4, font_size=12, color=rgb(150,170,200), align=PP_ALIGN.CENTER)
    add_rect(sl, 0, 7.38, 4.44, 0.12, fill_rgb=SAFFRON)
    add_rect(sl, 4.44, 7.38, 4.44, 0.12, fill_rgb=WHITE)
    add_rect(sl, 8.88, 7.38, 4.45, 0.12, fill_rgb=GREEN)

    prs.save(r"C:\PropertyLegalApp\DigiSampatti_Documents\DigiSampatti_Minister_MP.pptx")
    print("Minister/MP deck saved.")


# ─────────────────────────────────────────────
# VC TEASER DECK
# ─────────────────────────────────────────────
def build_vc_deck():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    BLACK  = rgb(10, 10, 10)
    DKGRAY = rgb(25, 25, 45)
    GREEN  = rgb(0, 200, 83)
    WHITE  = rgb(255, 255, 255)
    LGRAY  = rgb(248, 248, 248)
    MGRAY  = rgb(224, 224, 224)
    GOLD   = rgb(249, 168, 37)

    blank = prs.slide_layouts[6]

    # ── SLIDE 1: COVER ────────────────────────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, BLACK)
    add_rect(sl, 0, 0, 13.33, 0.06, fill_rgb=rgb(0,180,70))
    add_text(sl, "DigiSampatti", 0, 1.3, 13.33, 1.3, font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, "ಡಿಜಿ ಸಂಪತ್ತಿ", 0, 2.5, 13.33, 0.65, font_size=24, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(sl, "India's first full-journey property intelligence platform.", 0, 3.1, 13.33, 0.5, font_size=20, color=rgb(200,210,220), align=PP_ALIGN.CENTER)
    add_text(sl, "Camera. GPS. Expert. Verdict. From advance payment to registration.", 0, 3.55, 13.33, 0.5, font_size=16, color=rgb(150,165,180), align=PP_ALIGN.CENTER)
    add_text(sl, "SEED ROUND  ·  INVESTOR TEASER  ·  CONFIDENTIAL", 0, 4.25, 13.33, 0.4, font_size=12, color=rgb(100,115,130), align=PP_ALIGN.CENTER)
    for i, badge in enumerate(["Patent Pending", "Karnataka · India", "Google Play Store (Submitted)", "DPDP Act Compliant"]):
        add_rect(sl, 0.8 + i*3.0, 5.1, 2.7, 0.42, fill_rgb=rgb(25,25,40), line_rgb=rgb(60,80,100), line_width=1)
        add_text(sl, badge, 0.82 + i*3.0, 5.13, 2.66, 0.36, font_size=11, bold=True, color=rgb(180,195,210), align=PP_ALIGN.CENTER)
    add_rect(sl, 0, 7.44, 13.33, 0.06, fill_rgb=GREEN)

    # ── SLIDE 2: THE PROBLEM ──────────────────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, WHITE)
    add_rect(sl, 0, 0, 13.33, 0.08, fill_rgb=BLACK)
    add_text(sl, "The Problem", 0.5, 0.15, 12, 0.38, font_size=11, bold=True, color=GREEN)
    add_text(sl, "₹50,000 Crore Problem. Zero Platforms Solving It Completely.", 0.5, 0.5, 12.33, 0.75, font_size=27, bold=True, color=BLACK)

    add_rect(sl, 0.5, 1.3, 12.33, 0.72, fill_rgb=BLACK)
    add_text(sl, "India: 6 crore property transactions/year.  90% of buyers verify nothing.  Average lawyer fee: ₹50,000.  DigiSampatti: ₹99.", 0.65, 1.35, 12, 0.62, font_size=14, bold=True, color=WHITE)

    probs = [
        ("7 Portals. No Connection.", "Bhoomi, IGRS, RERA, BBMP, eCourts, CERSAI, Benami — each accurate in isolation. Fraud only appears when all 7 are read against each other. No platform does this."),
        ("Advance Paid Blind.", "Buyer pays ₹1–5L token money before due diligence. Discovers issues after paying. Fights months for refund. No pre-advance check tool exists."),
        ("Journey Broken After Step 1.", "Even if a buyer runs one check, they have no help with legal opinion, stamp duty, registration prep, mutation tracking, or portfolio management."),
    ]
    for i, (title, desc) in enumerate(probs):
        x = 0.5 + i*4.2
        add_rect(sl, x, 2.15, 4.0, 2.5, fill_rgb=BLACK, line_rgb=GREEN, line_width=1)
        add_text(sl, title, x+0.15, 2.22, 3.7, 0.45, font_size=13, bold=True, color=GREEN)
        add_text(sl, desc, x+0.15, 2.65, 3.7, 1.85, font_size=12, color=rgb(200,210,220))

    metrics = [("6Cr", "transactions/year\nin India"), ("90%", "buyers verify\nnothing"), ("₹50K", "avg lawyer fee\nper check"), ("₹99", "DigiSampatti\ncomplete check")]
    for i, (num, lbl) in enumerate(metrics):
        x = 0.5 + i*3.2
        add_rect(sl, x, 4.85, 3.0, 1.5, fill_rgb=LGRAY, line_rgb=MGRAY, line_width=1)
        col = GREEN if num == "₹99" else BLACK
        add_text(sl, num, x, 4.9, 3.0, 0.75, font_size=34, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(sl, lbl, x, 5.6, 3.0, 0.65, font_size=11, color=rgb(120,120,120), align=PP_ALIGN.CENTER)

    # ── SLIDE 3: FULL JOURNEY ─────────────────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, WHITE)
    add_rect(sl, 0, 0, 13.33, 0.08, fill_rgb=BLACK)
    add_text(sl, "The Differentiator", 0.5, 0.15, 12, 0.38, font_size=11, bold=True, color=GREEN)
    add_text(sl, "Landeed Ends Where DigiSampatti Begins", 0.5, 0.5, 12.33, 0.7, font_size=28, bold=True, color=BLACK)
    add_text(sl, "Landeed: step 4 of a 10-step journey. DigiSampatti: all 10 steps.", 0.5, 1.15, 12, 0.38, font_size=14, italic=True, color=rgb(100,100,100))

    headers = ["Journey Step", "Today (Without DigiSampatti)", "DigiSampatti Solution"]
    data = [
        ["Before paying advance", "Pays ₹1–5L token money blind. Discovers issues after.", "Full 7-portal check before a single rupee changes hands."],
        ["Property identification", "Doesn't know survey number. No digital entry point.", "Camera scans document. GPS identifies land on-site. Auto-extract."],
        ["Full verification", "3–5 days. 7 portals. ₹15,000 lawyer. No verdict.", "10 min. ₹99. 7 portals. Kannada verdict. Govt fees included."],
        ["Expert legal opinion", "₹50,000. Days of wait. No rural access.", "Verified Karnataka lawyer. ₹500–₹2,000. 2 hours. Remote."],
        ["Stamp duty calculation", "Surprise at sub-registrar office.", "Exact calculation by type, value, district — upfront."],
        ["Mutation tracking", "Manual follow-up for months. No notification.", "Auto-SMS when Bhoomi mutation confirmed."],
    ]
    add_table(sl, data, headers, 0.5, 1.6, 12.33, 5.2,
              hdr_bg=BLACK, hdr_fg=WHITE, row_bg1=WHITE, row_bg2=rgb(245,250,245))

    # ── SLIDE 4: THREE INPUTS + GOVT FEES ─────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, WHITE)
    add_rect(sl, 0, 0, 13.33, 0.08, fill_rgb=BLACK)
    add_text(sl, "The Product", 0.5, 0.15, 12, 0.38, font_size=11, bold=True, color=GREEN)
    add_text(sl, "Three Ways In. One Complete Answer.", 0.5, 0.5, 12.33, 0.7, font_size=28, bold=True, color=BLACK)

    inputs = [
        ("📷  Camera", "Photograph any property document — RTC, sale deed, mutation certificate. OCR extracts owner name, survey number, village, taluk, district. User confirms. Full check runs. Works for farmers who cannot type."),
        ("📍  GPS — Stand on Land", "Stand on the property. Tap once. GPS coordinates query Karnataka's Bhoomi cadastral GIS layer. Survey number identified automatically. 7-portal verification begins. No document needed."),
        ("⌨️  Type Survey Number", "Fastest path for buyers who know property ID. One entry, full report. Used by urban buyers, NRIs verifying remotely, banks processing loans, and real estate professionals."),
    ]
    for i, (title, desc) in enumerate(inputs):
        x = 0.5 + i*4.2
        add_rect(sl, x, 1.3, 4.0, 2.8, fill_rgb=rgb(232,245,233), line_rgb=rgb(165,214,167), line_width=1)
        add_text(sl, title, x+0.15, 1.37, 3.7, 0.48, font_size=15, bold=True, color=rgb(27,70,0))
        add_text(sl, desc, x+0.15, 1.83, 3.7, 2.15, font_size=12, color=rgb(46,125,50))

    add_rect(sl, 0.5, 4.28, 12.33, 0.85, fill_rgb=rgb(255,253,231), line_rgb=GOLD, line_width=2)
    add_text(sl, "Government Portal Fees — Fully Paid:  Bhoomi ₹15–25  ·  IGRS ₹30–50  ·  CERSAI ₹10  ·  BBMP ₹100–500\nAll included in the ₹99 service fee. Government revenue is maintained and increases — citizens who never checked portals individually now pay those fees through DigiSampatti.", 0.65, 4.3, 12, 0.81, font_size=12, bold=True, color=rgb(78,40,0))

    add_rect(sl, 0.5, 5.28, 12.33, 0.72, fill_rgb=LGRAY, line_rgb=MGRAY, line_width=1)
    add_text(sl, "Full technical architecture — how 7 portals are queried, how contradictions are detected, how the Karnataka legal rule engine works, and how the Kannada verdict is generated", 0.65, 5.3, 12, 0.68, font_size=12, italic=True, color=rgb(160,160,160))
    add_text(sl, "Available under NDA.", 0.65, 5.9, 12, 0.35, font_size=13, bold=True, color=rgb(100,100,100))

    # ── SLIDE 5: EXPERT COUNCIL + AI TRAINING ─
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, WHITE)
    add_rect(sl, 0, 0, 13.33, 0.08, fill_rgb=BLACK)
    add_text(sl, "The Intelligence Moat", 0.5, 0.15, 12, 0.38, font_size=11, bold=True, color=GREEN)
    add_text(sl, "AI Trained by Domain Experts — The Way the Best AI Is Built", 0.5, 0.5, 12.33, 0.7, font_size=27, bold=True, color=BLACK)
    add_text(sl, "Anthropic hired 20 financial analysts to train Claude on finance. DigiSampatti built the same model for Karnataka property law.", 0.5, 1.15, 12.33, 0.42, font_size=13, italic=True, color=rgb(80,80,80))

    council = [
        ("🏛  Retired Revenue Inspectors (5)", "Former Karnataka Revenue Dept officers. Know Bhoomi's mutation registers, Column 11/12 meanings, PTCL Act, taluk-specific quirks no tech team can learn from docs alone. Review borderline cases, correct the AI."),
        ("⚖️  Bar Council Lawyers (8)", "Verified property lawyers reviewing CAUTION cases. Apply Karnataka Land Reform Act nuances. Flag where rule engine misses regional legal distinctions. Train AI on real-world disputes — not textbook law."),
        ("🏦  Bank Loan Officers — Empanelled (5)", "Senior officers from Karnataka's public sector banks. Review CERSAI and IGRS interpretations. Train AI on actual mortgage risk vs cleared charge — institutional knowledge no public dataset contains."),
        ("🔏  Retired Sub-registrar Officers (2)", "Know the registration process, common documentation errors, red flags in sale deeds. Train AI on registration-stage patterns that would otherwise be missed entirely."),
    ]
    for i, (title, desc) in enumerate(council):
        y = 1.68 + i*1.2
        add_rect(sl, 0.5, y, 0.08, 0.95, fill_rgb=GREEN)
        add_rect(sl, 0.68, y, 12.15, 0.95, fill_rgb=LGRAY, line_rgb=MGRAY, line_width=1)
        add_text(sl, title, 0.82, y+0.06, 11.8, 0.38, font_size=13, bold=True, color=BLACK)
        add_text(sl, desc, 0.82, y+0.45, 11.8, 0.44, font_size=11, color=rgb(80,80,80))

    add_rect(sl, 0.5, 6.5, 12.33, 0.72, fill_rgb=BLACK)
    add_text(sl, "After 1,000 checks → AI recognises Karnataka fraud patterns.  After 10,000 → AI predicts high-risk taluks.  After 1,00,000 → Most accurate property intelligence dataset in India. Owned by DigiSampatti.", 0.65, 6.53, 12, 0.66, font_size=12, bold=True, color=WHITE)

    # ── SLIDE 6: vs LANDEED ───────────────────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, WHITE)
    add_rect(sl, 0, 0, 13.33, 0.08, fill_rgb=BLACK)
    add_text(sl, "Competitive Position", 0.5, 0.15, 12, 0.38, font_size=11, bold=True, color=GREEN)
    add_text(sl, "DigiSampatti vs Landeed ($16.3M Raised)", 0.5, 0.5, 12.33, 0.7, font_size=28, bold=True, color=BLACK)

    headers = ["Feature", "Landeed", "DigiSampatti"]
    data = [
        ["Core product", "Document download service", "Full-journey property intelligence platform"],
        ["Output", "PDF documents", "SAFE / CAUTION / DO NOT BUY verdict in Kannada"],
        ["Input method", "Type survey number only", "Camera scan / GPS on land / type"],
        ["Portals checked", "2–3 portals", "7 portals incl. Benami + eCourts reverse lookup"],
        ["Cross-portal contradiction", "No", "Yes — fraud only visible across portals"],
        ["Pre-advance check", "No", "Yes — before paying ₹1L token money"],
        ["Expert lawyer network", "No", "Yes — verified Karnataka lawyers, ₹500–₹2,000"],
        ["Stamp duty calculator", "No", "Yes — by property type, value, district"],
        ["Mutation tracker", "No", "Yes — auto-SMS when Bhoomi mutation complete"],
        ["AI trained by domain experts", "No", "Yes — Expert Council of 20 Karnataka professionals"],
        ["Language", "English only", "Kannada + English — accessible to every citizen"],
        ["Price per complete check", "₹299–499 per document", "₹99 — all 7 portals, all fees, complete report"],
        ["Patent protection", "None on core methodology", "Patent pending — cross-portal intelligence method"],
    ]
    add_table(sl, data, headers, 0.5, 1.28, 12.33, 5.5,
              hdr_bg=BLACK, hdr_fg=WHITE, row_bg1=WHITE, row_bg2=LGRAY)

    # ── SLIDE 7: MARKET ───────────────────────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, WHITE)
    add_rect(sl, 0, 0, 13.33, 0.08, fill_rgb=BLACK)
    add_text(sl, "Market", 0.5, 0.15, 12, 0.38, font_size=11, bold=True, color=GREEN)
    add_text(sl, "A Market That Has Never Been Fully Served", 0.5, 0.5, 12.33, 0.7, font_size=28, bold=True, color=BLACK)

    metrics2 = [("₹300Cr", "Karnataka\nYear 1–2"), ("₹8,500Cr", "India TAM\nproperty verification"), ("$16.3M", "Landeed raised\n1 portal, 1 state"), ("7×", "More portals\nthan Landeed")]
    for i, (num, lbl) in enumerate(metrics2):
        x = 0.5 + i*3.2
        add_rect(sl, x, 1.3, 3.0, 1.35, fill_rgb=LGRAY, line_rgb=MGRAY, line_width=1)
        col = GREEN if num == "7×" else BLACK
        add_text(sl, num, x, 1.35, 3.0, 0.72, font_size=34, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(sl, lbl, x, 2.0, 3.0, 0.58, font_size=11, color=rgb(120,120,120), align=PP_ALIGN.CENTER)

    headers = ["Segment", "Size", "DigiSampatti's Entry"]
    data = [
        ["Direct Consumers — buyers, NRIs, farmers", "1.2Cr Karnataka transactions/year", "₹99/check. Camera + GPS. Kannada. Rural reach."],
        ["Banks and NBFCs — pre-sanction checks", "45+ Karnataka institutions", "B2B API. Bulk verification. NPA prevention."],
        ["Real Estate Platforms — MagicBricks, 99acres", "Millions of listings, no verification", "DigiSampatti Verified badge on listings."],
        ["Expert Network — lawyers, revenue consultants", "15,000+ Karnataka property lawyers", "Referral marketplace. ₹500–₹2,000 per review."],
        ["Government / PMAY beneficiaries", "Lakhs of beneficiaries per year", "Pre-disbursement verification. Govt partnership."],
    ]
    add_table(sl, data, headers, 0.5, 2.75, 12.33, 3.85, font_size=11)

    add_rect(sl, 0.5, 6.72, 12.33, 0.55, fill_rgb=LGRAY, line_rgb=MGRAY, line_width=1)
    add_text(sl, "Unit economics, LTV/CAC, revenue model, state expansion roadmap — Available under NDA.", 0.65, 6.75, 12, 0.5, font_size=12, italic=True, color=rgb(140,140,140))

    # ── SLIDE 8: TRACTION ─────────────────────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, WHITE)
    add_rect(sl, 0, 0, 13.33, 0.08, fill_rgb=BLACK)
    add_text(sl, "Traction", 0.5, 0.15, 12, 0.38, font_size=11, bold=True, color=GREEN)
    add_text(sl, "Built. Filed. Submitted. In Conversation.", 0.5, 0.5, 12.33, 0.7, font_size=28, bold=True, color=BLACK)

    cols_data = [
        ("Product Status", [
            "✓ Flutter app — submitted to Google Play Store",
            "✓ 7-portal verification engine built",
            "✓ Camera OCR document reader — in development",
            "✓ GPS + Bhoomi GIS integration — in development",
            "✓ Kannada language interface complete",
            "✓ Expert Council structure defined",
            "✓ Stamp duty calculator — in development",
            "✓ Mutation tracker — in development",
        ]),
        ("IP, Legal, and Market", [
            "✓ Patent filed — core methodology",
            "✓ DPDP Act 2023 compliant architecture",
            "✓ Government fees payment model designed",
            "✓ Revenue Minister presentation prepared",
            "✓ MP constituency pilot discussions live",
            "✓ API access formally requested",
            "✓ Influencer partnerships in discussion",
            "✓ NRI community (UAE/UK) interest confirmed",
        ]),
    ]
    for i, (title, items) in enumerate(cols_data):
        x = 0.5 + i*6.45
        add_rect(sl, x, 1.3, 6.2, 5.3, fill_rgb=BLACK)
        add_text(sl, title, x+0.15, 1.37, 5.9, 0.45, font_size=14, bold=True, color=GREEN)
        add_text(sl, "\n".join(items), x+0.15, 1.85, 5.9, 4.6, font_size=12, color=rgb(200,210,220))

    # ── SLIDE 9: FOUNDER ──────────────────────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, WHITE)
    add_rect(sl, 0, 0, 13.33, 0.08, fill_rgb=BLACK)
    add_text(sl, "Founder", 0.5, 0.15, 12, 0.38, font_size=11, bold=True, color=GREEN)
    add_text(sl, "Why This Founder. Why Karnataka. Why Now.", 0.5, 0.5, 12.33, 0.7, font_size=28, bold=True, color=BLACK)

    add_rect(sl, 0.5, 1.3, 12.33, 0.72, fill_rgb=BLACK)
    add_text(sl, "Sanjay R — Karnataka-native. Technology + cybersecurity background.\nBuilt DigiSampatti end-to-end: product, legal architecture, Expert Council, government strategy, IP filing.", 0.65, 1.33, 12, 0.66, font_size=13, bold=True, color=WHITE)

    adv_text = ("Founder Advantages:\n\n"
                "→ Karnataka-native — deep knowledge of Bhoomi, taluk systems, local law\n"
                "→ Cybersecurity background — DPDP-compliant architecture from day one\n"
                "→ Solo founder — speed advantage, no consensus delay\n"
                "→ Government-first approach — partner model, not disruption\n"
                "→ Patent filed before investor conversations — IP protected before dilution\n"
                "→ Full-journey vision from day one — not a document tool that expanded")
    add_rect(sl, 0.5, 2.15, 6.0, 4.55, fill_rgb=LGRAY, line_rgb=MGRAY, line_width=1)
    add_text(sl, adv_text, 0.65, 2.22, 5.7, 4.4, font_size=12, color=rgb(50,50,50))

    timing_text = ("Market Timing:\n\n"
                   "→ Bhoomi 8.0 + Kaveri 3.0 launched — best govt data ever\n"
                   "→ DPDP Act 2023 creates compliance framework proptech needs\n"
                   "→ Smartphone penetration in rural Karnataka at historic high\n"
                   "→ NRI property investment into Karnataka growing post-COVID\n"
                   "→ No funded competitor with full-journey model in Karnataka\n"
                   "→ Landeed raised $16.3M proving the market — DigiSampatti builds what they couldn't")
    add_rect(sl, 6.83, 2.15, 6.0, 4.55, fill_rgb=LGRAY, line_rgb=MGRAY, line_width=1)
    add_text(sl, timing_text, 6.98, 2.22, 5.7, 4.4, font_size=12, color=rgb(50,50,50))

    # ── SLIDE 10: THE ASK ─────────────────────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, WHITE)
    add_rect(sl, 0, 0, 13.33, 0.08, fill_rgb=BLACK)
    add_text(sl, "The Ask", 0.5, 0.15, 12, 0.38, font_size=11, bold=True, color=GREEN)
    add_text(sl, "Seed Round — Building India's Property Intelligence Layer", 0.5, 0.5, 12.33, 0.7, font_size=26, bold=True, color=BLACK)

    add_rect(sl, 0.5, 1.3, 12.33, 0.72, fill_rgb=BLACK)
    add_text(sl, "Stage: Pre-revenue, post-patent, app submitted   ·   Use: Product completion · Expert Council · Karnataka launch · Govt API", 0.65, 1.33, 12, 0.66, font_size=13, bold=True, color=WHITE)

    use_cols = [
        ("Product Completion", "Camera OCR engine. GPS + GIS integration. eCourts reverse lookup. Benami portal. Cross-portal contradiction engine. Stamp duty calculator. Mutation tracker. Kannada NLP for plain-language verdicts."),
        ("Expert Council", "Onboard 20 Karnataka domain experts — Revenue Inspectors, Bar Council lawyers, bank officers, sub-registrar officers. Build review workflow and AI training feedback loop. This is the accuracy moat."),
        ("Karnataka Market Launch", "5-district simultaneous launch. Sub-registrar office activations. Real estate agent network. NRI community outreach. Bank B2B conversations. Government API fast-track support."),
    ]
    for i, (title, desc) in enumerate(use_cols):
        x = 0.5 + i*4.2
        add_rect(sl, x, 2.15, 4.0, 2.7, fill_rgb=LGRAY, line_rgb=MGRAY, line_width=1)
        add_text(sl, title, x+0.12, 2.22, 3.76, 0.45, font_size=13, bold=True, color=BLACK)
        add_text(sl, desc, x+0.12, 2.68, 3.76, 2.05, font_size=12, color=rgb(70,70,70))

    add_rect(sl, 0.5, 5.0, 12.33, 0.72, fill_rgb=LGRAY, line_rgb=MGRAY, line_width=1)
    add_text(sl, "Exact raise amount, valuation, equity structure, financial projections, unit economics, 24-month revenue roadmap — Available under NDA. Schedule a call to receive full deck.", 0.65, 5.03, 12, 0.66, font_size=12, italic=True, color=rgb(140,140,140))

    # ── SLIDE 11: NDA CLOSING ─────────────────
    sl = prs.slides.add_slide(blank)
    slide_bg(sl, BLACK)
    add_rect(sl, 0, 0, 13.33, 0.06, fill_rgb=GREEN)
    add_text(sl, "You Have Seen the Shape. Not the Engine.", 0, 1.0, 13.33, 1.0, font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    reveal = ("What you have not seen yet:\n\n"
              "The proprietary cross-portal intelligence engine.\n"
              "The unit economics and LTV/CAC model.\n"
              "The government partnership strategy and API roadmap.\n"
              "The Expert Council activation plan.\n"
              "The exact raise terms and use-of-funds breakdown.\n"
              "The state expansion playbook beyond Karnataka.\n\n"
              "These details are withheld because they are the moat.\nNot because they are weak.")
    add_text(sl, reveal, 2.0, 1.95, 9.33, 3.5, font_size=14, color=rgb(190,200,215), align=PP_ALIGN.CENTER)

    add_rect(sl, 3.0, 5.2, 7.33, 1.55, fill_rgb=rgb(0,40,20), line_rgb=GREEN, line_width=2)
    add_text(sl, "Next Step", 3.0, 5.25, 7.33, 0.42, font_size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(sl, "Sign standard mutual NDA.\nReceive full technical deck, financial model, unit economics,\ngovernment partnership documents, and Expert Council plan.\n\nOne conversation. Complete transparency. No games.", 3.0, 5.65, 7.33, 1.0, font_size=12, color=rgb(200,215,200), align=PP_ALIGN.CENTER)

    add_text(sl, "Sanjay R  ·  rajsanjay381@gmail.com  ·  +91 89043 42255", 0, 6.95, 13.33, 0.38, font_size=13, color=GREEN, align=PP_ALIGN.CENTER)
    add_rect(sl, 0, 7.44, 13.33, 0.06, fill_rgb=GREEN)

    prs.save(r"C:\PropertyLegalApp\DigiSampatti_Documents\DigiSampatti_VC_Teaser.pptx")
    print("VC Teaser deck saved.")


build_minister_deck()
build_vc_deck()
print("Both PPTs generated successfully.")
